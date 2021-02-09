import os, time, uuid, sys
from urllib.parse import quote as urlquote
import asyncio
import datetime

from django.contrib import admin
from django.utils.html import mark_safe
from django.urls import reverse
from django.utils.html import format_html
from django import forms
from django.http import HttpResponseRedirect
from django.contrib import messages
from django.core.files import File
from django.utils.translation import gettext as _, ngettext
from django.conf import settings
from asgiref.sync import sync_to_async
from asgiref.sync import async_to_sync
from django.utils.crypto import get_random_string
from django.shortcuts import get_object_or_404, redirect, render

from sorl.thumbnail.admin import AdminImageMixin
from PIL import Image
from crispy_forms.helper import FormHelper
from crispy_forms.layout import Submit

from .models import *

# references:
# for sorl.thumbnail: https://sorl-thumbnail.readthedocs.io/en/latest/examples.html
############################################################################################################
############################################################################################################
############################################################################################################
class BipageInline(admin.TabularInline):
    model = Bipage
    show_change_link = True
    fields = ['image_slidepage_preview', 'name', 'ppt_page_layout', 'last_updated', 'edit']
    readonly_fields = ['last_updated', 'image_slidepage_preview', 'edit']
    list_display = ['image_slidepage_preview', 'name', 'ppt_page_layout', 'last_updated']
    can_delete = True

    # prepopulated_fields = {"name": ("",)}
    extra = 1

    list_per_page = 5

    per_page = 5

    def image_slidepage_preview(self, obj):
        return mark_safe("<img src='/static/biportal/img/SlidePage.png' />")

    def edit(self, instance):
        # url = reverse('admin:bipage'
        # return format_html(u'<a href="{}">Edit</a>', url)
        # # … or if you want to include other fields:
        # return format_html(u'<a href="{}">Edit: {}</a>', '/admin/bipage', instance.title)
        html = format_html("""
        <a href="/admin/biportal/bipage/{}/change/">Edit</a>
        """, instance.pk)

        return html
############################################################################################################
############################################################################################################
############################################################################################################
class PresentatioAdminForm(forms.ModelForm):
    def __init__(self, *args, **kwargs):
        super(PresentatioAdminForm, self).__init__(*args, **kwargs)

        # self.fields['embedded'].label = 'Embedded report'
        self.fields['description'].widget = forms.Textarea(attrs={'cols': '40', 'rows': '10'})

    class Meta:
        model = Presentation
        fields = '__all__'
@admin.register(Presentation)
class PresentationAdmin(admin.ModelAdmin):

    form = PresentatioAdminForm

    save_on_top = True

    readonly_fields = ['created_at']

    fieldsets = [
        (
            'Basic Information',
            {
                'fields': ['name', 'description', 'ppt_master_file', 'active']
            }
        ),
        (
            'Metadata',
            {
                'fields': ['created_by', 'created_at', 'tags']
            }
        ),
        (
            'Exported files',
            {
                'classes': ('collapse',),
                'fields': [('ppt_file', 'pdf_file')]
            }
        ),
    ]

    search_fields = ['name', 'description',]

    list_filter = ['tags', 'created_by', 'active']

    list_display = ['name', 'created_by', 'created_at','updated_at', 'active']

    list_editable = ['active']

    inlines = [BipageInline]

    save_as = True

    list_per_page = 5

    actions = ['render_all_snippets_in_selected_presentations']

    def render_all_snippets_in_selected_presentations(self, request, queryset):
        pass
        # queryset.update(status='p')

    # change the context before the change form is rendered
    def render_change_form(self, request, context, *args, **kwargs):
        """We need to update the context to show the button."""

        # add a new key to the context
        context.update({'show_render_all_report_snippets': True})
        context.update({'show_export_ppt': True})
        context.update({'show_export_pdf': True})
        return super().render_change_form(request, context, *args, **kwargs)

    def response_post_save_change(self, request, obj):
        """This method is called by `self.changeform_view()` when the form
        was submitted successfully and should return an HttpResponse.
        """
        # Check that you clicked the button `_render_snippets`
        if '_render_snippets' in request.POST:

            # Create a copy of your object
            # Assuming you have a method `create_from_existing()` in your manager
            # new_obj = self.model.objects.create_from_existing(obj)
            # do something with obj

            # Get its admin url
            opts = self.model._meta
            info = self.admin_site, opts.app_label, opts.model_name
            route = '{}:{}_{}_change'.format(*info)

            # post_url = reverse(route, args=(new_obj.pk,))
            # post_url = reverse(route, args=(obj.pk,))
            post_url = "/admin/biportal/presentation/{}/change/".format(obj.pk)

            # And redirect
            return HttpResponseRedirect(post_url)

        else:

            # Otherwise, use default behavior
            return super().response_post_save_change(request, obj)
###################################################################################################
###################################################################################################
###################################################################################################
class SnippetAdminForm(forms.ModelForm):
    def __init__(self, *args, **kwargs):
        super(SnippetAdminForm, self).__init__(*args, **kwargs)

        self.fields['embedded'].widget = forms.Textarea(attrs={'cols': '80', 'rows': '3'})

    class Meta:
        model = Snippet
        fields = '__all__'
##################################################################################################
@admin.register(Snippet)
class SnippetAdmin(admin.ModelAdmin):

    form = SnippetAdminForm

    save_on_top = True

    change_form_template = 'admin/change_form_snippets.html'

    ordering = ('created_at',)

    readonly_fields = ['created_at','report_rendered_preview', 'report_snippet_preview', 'render_button']

    list_filter = ['tags', 'created_by', 'created_at']

    fieldsets = [
        (
            'Basic Information',
            {
                'fields': ['name', 'embedded', 'render_button']
            }
        ),
        (
            'Rendered report',
            {
                'classes': ('collapse',),
                'fields':
                [
                    ('report_rendered_preview', 'image_rendered'),
                ]
            }
        ),
        (
            'Snippet Image',
            {
                'fields':
                [
                    ('report_snippet_preview', 'image_cropped',)
                ]
            }
        ),
        (
            'Position and Dimensions',
            {
                'classes': ('collapse',),
                'fields': [('x', 'y'), ('w', 'h')],
                # 'description': 'sfdsfd dsf sdf ds fds  fds  f'
            }
        ),

        (
            'Metadata',
            {
                'fields': ['created_by', 'created_at', 'tags']
            }
        ),
    ]

    list_display = ['image_cropped_preview', 'name', 'created_by', 'created_at','updated_at']

    list_display_links = ['image_cropped_preview', 'name']

    search_fields = ['name', 'created_at']

    save_as = True

    list_per_page = 10

    # this is adding an extra action to the list of actions in the list form
    actions = ['render_all_reports_in_selected_snippets']

    # reference:
    # https://ilovedjango.com/django/admin/how-to-show-image-from-imagefield-in-django-admin-page/
    def report_rendered_preview(self, obj):
        # return mark_safe("""<img src={url}
        # width={width}
        # height={height} id='image'/>""".format(
        #     url = obj.image_rendered.url,
        #     width=obj.image_rendered.width,
        #     height=obj.image_rendered.height
        #     ))
        return mark_safe("""
    <div class="modal fade" id="modalCrop">
        <div class="modal-dialog">
            <div class="modal-content">
                <div class="modal-footer">
                <!-- <button type="button" class="btn btn-default" data-dismiss="modal">Cancel</button> -->
                <button type="button" class="btn btn-primary js-crop-and-upload">Set Snippet Position</button>
                </div>
                <div class="modal-body">
                    <img src={url} width={width} height={height} id="image2" style="max-width: 100%;">
                </div>
                <div class="modal-footer">
                <!-- <button type="button" class="btn btn-default" data-dismiss="modal">Cancel</button> -->
                <button type="button" class="btn btn-primary js-crop-and-upload">Set Snippet Position</button>
                </div>
            </div>
        </div>
    </div>
        """.format(
            url = obj.image_rendered.url,
            width=obj.image_rendered.width,
            height=obj.image_rendered.height
            ))

    def report_snippet_preview(self, obj):

        html = ''
        url = ''
        try:
            url = obj.image_cropped.url
        except:
            pass

        html = """
        <div id='snippet_preview' style='overflow: hidden'>
            <img src={url} />
        </div>
        """.format(url = url,)

        return mark_safe(html)

    def image_cropped_preview(self, obj):

        try:
            html = mark_safe("<img src={url} width=220 />".format(
                url = obj.image_cropped.url,
                ))
        except:
            html = ""

        return html

    def image_rendered_preview(self, obj):
        return mark_safe("<img src={url} width={width} height={height} />".format(
            url = obj.image_rendered.url,
            width=obj.image_rendered.width,
            height=obj.image_rendered.height
            ))

    ##########################################################################################
    # default created_by user with curent user
    # overwrite get_changeform_initial_data
    # called before the add form is rendered
    def get_changeform_initial_data(self, request):
        get_data = super(SnippetAdmin, self).get_changeform_initial_data(request)
        get_data['created_by'] = request.user.pk
        return get_data

    # augment the request object before rendering add form
    # reference: https://docs.djangoproject.com/en/dev/ref/contrib/admin/#django.contrib.admin.ModelAdmin.change_view
    def add_view(self, request, form_url='', extra_context=None):


        # set the default position and dimension of the cropping selection
        extra_context = {
            'xx': 100,
            'yy': 50,
            'hh': 240,
            'ww': 240,
        }

        return super().add_view(
            request, form_url, extra_context=extra_context
        )

    # augment the request object before rendering change form
    # reference: https://docs.djangoproject.com/en/dev/ref/contrib/admin/#django.contrib.admin.ModelAdmin.change_view
    def change_view(self, request, object_id, form_url='', extra_context=None):
        # extra_context = extra_context or {}
        # extra_context['osm_data'] = self.get_osm_info()

        # obj = Snippet.all
        obj = get_object_or_404(Snippet, pk=object_id)

        extra_context = {
            'xx': obj.x,
            'yy': obj.y,
            'hh': obj.h,
            'ww': obj.w,
        }

        return super().change_view(
            request, object_id, form_url, extra_context=extra_context,
        )



    # change the context before the change form is rendered
    def render_change_form(self, request, context, *args, **kwargs):
        """We need to update the context to show the button."""

        # add a new key to the context
        return super().render_change_form(request, context, *args, **kwargs)

    # code behind the extra action
    def render_all_reports_in_selected_snippets(self, request, queryset):
        pass
        # queryset.update(status='p')

    # handle button sumbissions
    # overwrite default behavior if you want to change the message
    # after submission
    # obj is the data
    def response_change(self, request, obj):

        if '_continue' in request.POST or '_save' in request.POST:
            try:
                if obj.x > 0.0 and obj.y > 0.0:
                    x = obj.x
                    y = obj.y
                    h = obj.h
                    w = obj.w
                    image = Image.open(obj.image_rendered.file)
                    cropped_image = image.crop((x, y, w+x, h+y))

                    uid = get_random_string(length=16, allowed_chars=u'abcdefghijklmnopqrstuvwxyzABCDEFGHIJKLMNOPQRSTUVWXYZ0123456789')
                    date = datetime.datetime.now()
                    result = '%s-%s-%s-%s-%s-%s_%s' % (date.year, date.month, date.day, date.hour, date.minute, date.second, uid)
                    filepath = "{}/image_cropped/{}.png".format(
                        settings.MEDIA_ROOT,
                        result
                    )

                    file_in_media = "image_cropped/{}.png".format(
                        result
                    )

                    cropped_image.save(filepath, format='PNG')

                    obj.image_cropped = file_in_media

                    obj.save()
            except:
                pass

        if '_render_report' in request.POST:
            # this preventing creating the default message for save
            return self.response_post_save_change(request, obj)

        return super().response_change(request, obj)

    # handle button submissions too
    # called after response_change
    # obj is the data
    def response_post_save_change(self, request, obj):
        """This method is called by `self.changeform_view()` when the form
        was submitted successfully and should return an HttpResponse.
        """
        # Check that you clicked the button `_render_snippets`
        if '_render_report' in request.POST:

            # obj is the model data
            # do something with obj

            # reference: saving image file to django ImageField
            # https://stackoverflow.com/questions/1308386/programmatically-saving-image-to-django-imagefield
            # https://stackoverflow.com/questions/13393191/programmatically-add-file-to-django-imagefield

            # obj.image_rendered = 'image_rendered/tableau-scrap-screenshot.png'
            # obj.save()

            #################################################################
            # logic to render the report using pyppeteer:

            # url = 'https://public.tableau.com/en-us/gallery/holiday-consumer-spending?tab=featured&type=featured'

            url = ''

            if obj.embedded:
                url = obj.embedded

            # # asyncio.set_event_loop_policy(asyncio.WindowsProactorEventLoopPolicy())
            # loop = asyncio.new_event_loop()
            # asyncio.set_event_loop(loop)
            # asyncio.get_event_loop().run_until_complete(render_report(url, filepath))

            rendered_ok = False
            file_in_media = ''

            try:
                if url:
                    # generate the filename
                    uid = get_random_string(length=16, allowed_chars=u'abcdefghijklmnopqrstuvwxyzABCDEFGHIJKLMNOPQRSTUVWXYZ0123456789')
                    date = datetime.datetime.now()
                    result = '%s-%s-%s-%s-%s-%s_%s' % (date.year, date.month, date.day, date.hour, date.minute, date.second, uid)
                    filepath = "{}/image_rendered/{}.png".format(
                        settings.MEDIA_ROOT,
                        result
                    )

                    file_in_media = "image_rendered/{}.png".format(
                        result
                    )

                    loop = asyncio.new_event_loop()
                    asyncio.set_event_loop(loop)
                    loop.run_until_complete(render_report(url, filepath))
                    rendered_ok = True
                else:
                    rendered_ok = False
            except:
                # sys.exc_info()[0]
                raise

            # # reference:
            # # https://docs.djangoproject.com/en/3.1/topics/async/
            # sync_get_data = async_to_sync(render_report_02(), force_new_loop=True)
            # time.sleep(5)

            # asyncio.run(render_report_02())
            # async_to_sync(render_report_02(), force_new_loop=True)()

            # await render_report_02()
            # time.sleep(6)

            # save the image in image_rendered field
            # obj.image_rendered = 'image_rendered/test_report_render_001.png'

            if rendered_ok:
                obj.image_rendered = file_in_media
                obj.save()

                opts = self.model._meta
                preserved_filters = self.get_preserved_filters(request)

                msg = _('The report was rendered successsfuly')
                self.message_user(request, msg, messages.SUCCESS)


            # post_url = reverse(route, args=(new_obj.pk,))
            # post_url = reverse(route, args=(obj.pk,))
            post_url = "/admin/biportal/snippet/{}/change/".format(obj.pk)

            # And redirect
            # return super().response_post_save_change(request, obj)
            return HttpResponseRedirect(post_url)

        # elif '_continue' in request.POST or '_save' in request.POST:
        #     x = obj.x
        #     y = obj.y
        #     h = obj.h
        #     w = obj.w
        #     image = Image.open(obj.image_rendered.file)
        #     cropped_image = image.crop((x, y, w+x, h+y))

        #     obj.image_cropped = cropped_image.file
        #     return super().response_post_save_change(request, obj)

        else:

            # Otherwise, use default behavior
            return super().response_post_save_change(request, obj)

        # reference:
        # https://docs.djangoproject.com/en/3.1/topics/forms/modelforms/
        # https://stackoverflow.com/questions/4670783/make-the-user-in-a-model-default-to-the-current-user
        # every time you save a form using commit=False, Django adds a save_m2m()
        # method to your ModelForm subclass. After you’ve manually saved the instance
        # produced by the form, you can invoke save_m2m()
        # to save the many-to-many form data. For example:
        def save_model(self, request, instance, form, change):
            user = request.user
            instance = form.save(commit=False)
            if not change or not instance.created_by:
                instance.created_by = user
            instance.updated_by = user

            ###########################
            # instance.name = 'test changed snippet 1001'

            instance.save()
            form.save_m2m()
            return instance

# admin.site.register(Presentation)
# admin.site.register(Snippet)
############################################################################################################
############################################################################################################
############################################################################################################
# TODO: should be moved to it's own class
from pyppeteer import launch
async def render_report(url, filepath):

    # browser = await launch()
    browser = await launch(
        handleSIGINT=False,
        handleSIGTERM=False,
        handleSIGHUP=False
        )
    page = await browser.newPage()

    await page.setViewport({
        'width': 1200,
        'height': 1000
    })

    await page.goto(url, {'waitUntil': 'networkidle2' })
    # time.sleep(5)
    await page.screenshot({
        'path': filepath,
        'fullPage': 'true'
        })

    dimensions = await page.evaluate('''() => {
        return {
            width: document.documentElement.clientWidth,
            height: document.documentElement.clientHeight,
            deviceScaleFactor: window.devicePixelRatio,
        }
    }''')

    print(dimensions)
    # >>> {'width': 800, 'height': 600, 'deviceScaleFactor': 1}
    await browser.close()
############################################################################################################
############################################################################################################
############################################################################################################
# admin.site.register(Bipage)
# admin.site.register(SnippetHtml)
class BipageAdminForm(forms.ModelForm):
    def __init__(self, *args, **kwargs):
        super(BipageAdminForm, self).__init__(*args, **kwargs)

        self.fields['name'].widget = forms.TextInput(attrs={'size': '28'})
        # self.fields['ppt_page_layout'].widget = forms.Select(attrs={'size': '40'})
        self.fields['title'].widget = forms.TextInput(attrs={'size': '70'})
        self.fields['subtitle'].widget = forms.TextInput(attrs={'size': '70'})
        # self.fields['ppt_page_layout'].widget = forms.Select(attrs={'size': '40'})

        # reference:
        # https://github.com/django-crispy-forms/django-crispy-forms/issues/697
        self.helper = FormHelper()
        self.helper.form_id = 'id-my-form'
        self.helper.form_class = 'my-form'
        self.helper.form_method = 'post'
        self.helper.add_input(Submit('submit', 'Submit'))

    class Meta:
        model = Bipage
        fields = '__all__'
#######################################################################################################
@admin.register(Bipage)
class BipageAdmin(admin.ModelAdmin):

    form = BipageAdminForm

    change_form_template = 'admin/change_form_bipage.html'

    save_on_top = True

    # save_as = True

    readonly_fields = ['presentation', 'layout_preview', 'page_preview']

    fieldsets = [
        (
            '',
            {
                'fields': ['presentation', 'name', 'ppt_page_layout', 'layout_preview']
            },
        ),
        (
            '',
            {
                'fields': [ ('title', 'subtitle',) ]
            },
        ),
        (
            '',
            {
                'fields': ['page_preview']
            }
        ),
        (
            'Components',
            {
                'classes': ('collapse',),
                'fields': [ ('snippets', 'texts'), ]
            },
        ),

        (
            'Exported ',
            {
                'classes': ('collapse',),
                'fields': [('ppt_file', 'pdf_file')]
            }
        ),
    ]

    # reference: hide it in the admin menu
    # https://stackoverflow.com/questions/2431727/django-admin-hide-a-model
    def get_model_perms(self, request):
        """
        Return empty perms dict thus hiding the model from admin index.
        """
        return {}

    def render_change_form(self, request, context, *args, **kwargs):
        """We need to update the context to show the button."""

        # add a new key to the context
        context.update({'show_render_all_report_snippets': True})
        context.update({'show_export_ppt': True})
        context.update({'show_export_pdf': True})
        context.update({'show_save_and_add_another': False})
        # context.update({'is_nav_sidebar_enabled': False})
        context.update({'show_presentation_nav': True})
        return super().render_change_form(request, context, *args, **kwargs)
        

    def response_post_save_change(self, request, obj):
        """ Custom actiaons for bi page
        Args:
            request ([type]): [description]
            obj ([type]): [description]
        Returns:
            [type]: [description]
        """
        if '_export_ppt' in request.POST:

            # logic to create a ppt out of the current bipage
            # TODO: move this to the model class, and invoke model method here

            # generate a filename for the ppt
            uid = get_random_string(length=16, allowed_chars=u'abcdefghijklmnopqrstuvwxyzABCDEFGHIJKLMNOPQRSTUVWXYZ0123456789')
            date = datetime.datetime.now()
            result = '%s-%s-%s-%s-%s-%s_%s' % (date.year, date.month, date.day, date.hour, date.minute, date.second, uid)
            filepath = "{}/generated_ppt/{}.pptx".format(
                settings.MEDIA_ROOT,
                result
            )
            file_in_media = 'generated_ppt/{}.pptx'.format(result)

            # core logic
            from pptx import Presentation
            from pptx.chart.data import CategoryChartData
            from pptx.enum.chart import XL_CHART_TYPE
            from pptx.chart.data import ChartData
            from pptx.util import Inches

            from biportal.templatetags.form_tags import replace_tokens

            master_file = 'ppt_master_files/guardian_master.pptx'

            # TODO: replace this with the master ppt file associated with the presentation
            master_file = f'{ settings.MEDIA_ROOT }/{ master_file }'

            # initialize the ppt presentation object
            prs = Presentation(master_file)

            # get the slide layout
            # this is based on the layout set for this bipage object

            # we default to the first layout
            this_slide_layout = prs.slide_layouts[0]
            if obj.ppt_page_layout == 'header':

                this_slide_layout = prs.slide_layouts[0]

                # add a slide to the presentation
                slide = prs.slides.add_slide(this_slide_layout)

                # get references to the title and subtitle
                title = slide.shapes.title
                subtitle = slide.placeholders[1]

                # set the values of title and subtitle to this bipage title and subtitle
                title.text = replace_tokens(obj.title)
                subtitle.text = replace_tokens(obj.subtitle)

            elif obj.ppt_page_layout == 'content1':

                this_slide_layout = prs.slide_layouts[1]
                # add a slide to the presentation
                slide = prs.slides.add_slide(this_slide_layout)

                # get references to the title and subtitle
                title = slide.shapes.title
                # subtitle = slide.placeholders[0]
                # subtitle = slide.shapes.subtitle

                # set the values of title and subtitle to this bipage title and subtitle
                title.text = replace_tokens(obj.title)
                # subtitle.text = replace_tokens(obj.subtitle)

                # for s in slide.shapes:
                #     s.text = s.name
                #     if s.name == 'Title 1':
                #         s.text = replace_tokens(obj.title)
                #     if s.name == 'Text Placeholder 3':
                #         s.text = replace_tokens(obj.subtitle)
                #     if s.name == 'Content Placeholder 2':
                #         # main placeholder
                #         # pass
                #         # s.text = replace_tokens(obj.subtitle)
                #         s.add_picture('/Users/marian.dumitrascu/Dropbox/Work/Current/python-cms/bi-portal/mainsite/media/image_cropped/2021-1-11_46OqBEruye4yfJD4.png', 0, 0)

                # for slideshape in slide.

                for shape in slide.placeholders:
                    shape.text = f'{shape.placeholder_format.idx}:{shape.name}'

                if obj.snippets.all():
                    picture_path = obj.snippets.all()[0].image_cropped.path

                    slide.shapes.add_picture(
                        picture_path,
                        Inches(0.5),
                        Inches(1.5),
                        height = Inches(5.0))

                slide.placeholders[0].text = replace_tokens(obj.title)
                slide.placeholders[13].text = replace_tokens(obj.subtitle)

                # slide.placeholders[1].text = obj.snippets.all()[0].image_cropped.url
                # slide.placeholders[1].add_picture('/Users/marian.dumitrascu/Dropbox/Work/Current/python-cms/bi-portal/mainsite/media/image_cropped/2021-1-11_46OqBEruye4yfJD4.png')
                # slide.shapes['Text Placeholder 3'].text =



            elif obj.ppt_page_layout == 'content2':

                this_slide_layout = prs.slide_layouts[2]
                # add a slide to the presentation
                slide = prs.slides.add_slide(this_slide_layout)

                # get references to the title and subtitle
                title = slide.shapes.title

                # set the values of title and subtitle to this bipage title and subtitle
                title.text = replace_tokens(obj.title)

                for shape in slide.placeholders:
                    shape.text = f'{shape.placeholder_format.idx}:{shape.name}'

                if obj.snippets.all():
                    picture_path = obj.snippets.all()[0].image_cropped.path

                    slide.shapes.add_picture(
                        picture_path,
                        Inches(0.5),
                        Inches(1.5),
                        height = Inches(5.0))

                slide.placeholders[0].text = replace_tokens(obj.title)
                slide.placeholders[15].text = replace_tokens(obj.subtitle)

            # save the presentation in media folder
            prs.save(filepath)

            # message the user with the link to the generated file
            # TODO: replace hardcoded base url (127.0.0.1)
            url_to_file = '/media/generated_ppt/{}.pptx'.format(result)
            msg = 'The PowerPoint for this page was generated successsfuly. You can downlad it from <a href={} target=_blank>here</a>'.format(url_to_file)
            msg = _(mark_safe(msg))
            self.message_user(request, msg, messages.SUCCESS)

            obj.ppt_file = file_in_media
            obj.save()

            # url to the edit page of this layout page
            post_url = "/admin/biportal/bipage/{}/change/".format(obj.pk)

            # redirect to the edit view of curent bipage
            return HttpResponseRedirect(post_url)

        else:

            return super().response_post_save_change(request, obj)

    def response_change(self, request, obj):
        if '_export_ppt' in request.POST:
            return self.response_post_save_change(request, obj)
        else:
            return super().response_change(request, obj)
############################################################################################################
############################################################################################################
############################################################################################################
class SnippetHtmlForm(forms.ModelForm):
    def __init__(self, *args, **kwargs):
        super(SnippetHtmlForm, self).__init__(*args, **kwargs)

        # self.fields['embedded'].label = 'Embedded report'
        self.fields['htmltext'].widget = forms.Textarea(attrs={'cols': '80', 'rows': '10'})

    class Meta:
        model = SnippetHtml
        fields = '__all__'
#####################################################################################################
@admin.register(SnippetHtml)
class SnippetHtmlAdmin(admin.ModelAdmin):

    form = SnippetHtmlForm

    readonly_fields = ['get_htmltext_as_markdown']

    fields = ['name', 'get_htmltext_as_markdown', 'htmltext']

    searchfields = ['name', 'htmltext']

    change_form_template = 'admin/change_form_markdown.html'

    # reference: hide it in the admin menu
    # https://stackoverflow.com/questions/2431727/django-admin-hide-a-model
    def get_model_perms(self, request):
        """
        Return empty perms dict thus hiding the model from admin index.
        """
        return {}


